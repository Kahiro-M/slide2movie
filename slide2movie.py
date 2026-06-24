import os
from pptx import Presentation
from gtts import gTTS
from pydub import AudioSegment
from pathlib import Path
import shutil
import argparse
import sys

# 標準出力をUTF-8に再設定（Windows環境での文字化け対策）
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

# stdout をファイルと元のストリームに同時に書き出すラッパー
class _Tee:
    def __init__(self, stream, filepath: str):
        # 元ストリームを行バッファリングに設定してからラップ
        if hasattr(stream, 'reconfigure'):
            # バイナリモード（rb/wb/ab）でなければ行バッファリングを設定
            mode = getattr(stream, 'mode', 'w')
            if mode not in ('rb', 'wb', 'ab'):
                try:
                    stream.reconfigure(line_buffering=True)
                except Exception:
                    pass
        self._stream = stream
        self._file = open(filepath, "w", encoding="utf-8", buffering=1)

    def write(self, data):
        self._stream.write(data)
        self._file.write(data)

    def flush(self):
        self._stream.flush()
        self._file.flush()

    def close(self):
        self._file.close()

    # sys.stdout が持つ属性に委譲（subprocess等の互換性のため）
    def __getattr__(self, name):
        return getattr(self._stream, name)
    
# スクリプトと同階層のffmpegを使用
_BASE_DIR = Path(__file__).parent

# ffmpeg / ffprobe のパスを解決する（クロスプラットフォーム対応）
def find_ffmpeg(name: str) -> str:
    import platform
    # システムの PATH 上にあれば優先
    found = shutil.which(name)
    if found:
        return found
    # Windows の場合は同階層の .exe にフォールバック
    if platform.system() == "Windows":
        local = _BASE_DIR / f"{name}.exe"
        if local.exists():
            return str(local)
    raise FileNotFoundError(f"{name} が見つかりません。インストールまたは同階層に配置してください。")

AudioSegment.converter = find_ffmpeg("ffmpeg")
AudioSegment.ffmpeg    = find_ffmpeg("ffmpeg")
AudioSegment.ffprobe   = find_ffmpeg("ffprobe")

# ──────────────────────────────────────────
# 1. PPTXの各スライドをPNG画像に変換
# ──────────────────────────────────────────
# PPTXの各スライドをPNG画像に変換します。
# LibreOfficeを使わず、python-pptx + Pillowで実装。
# Args:
#     pptx_path (str): 入力PPTXファイルパス
#     output_dir (str): 出力ディレクトリ
#     dpi (int): 画像解像度（デフォルト150）
# Returns:
#     list[str]: 生成されたPNGファイルパスのリスト（スライド順）
def pptx_to_pngs(pptx_path, output_dir="slides_png", dpi=150):
    from PIL import Image, ImageDraw, ImageFont
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)

    # スライドサイズ取得（EMU → ピクセル変換）
    emu_per_inch = 914400
    slide_width_px  = int(prs.slide_width  / emu_per_inch * dpi)
    slide_height_px = int(prs.slide_height / emu_per_inch * dpi)

    png_paths = []

    for i, slide in enumerate(prs.slides):
        # スライドの背景色を取得（取得できない場合は白）
        try:
            bg_color = slide.background.fill.fore_color.rgb
            bg = tuple(int(bg_color[j:j+2], 16) for j in (0, 2, 4))
        except Exception:
            bg = (255, 255, 255)

        img = Image.new("RGB", (slide_width_px, slide_height_px), color=bg)
        draw = ImageDraw.Draw(img)

        # テキストシェイプを描画（簡易レンダリング）
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # 位置をEMU → ピクセルに変換
                x = int(shape.left  / emu_per_inch * dpi) if shape.left  else 10
                y = int(shape.top   / emu_per_inch * dpi) if shape.top   else 10
                draw.text((x, y), text, fill=(0, 0, 0))

        output_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
        img.save(output_path, "PNG", dpi=(dpi, dpi))
        print(f"PNG保存: {output_path}", flush=True)
        png_paths.append(output_path)

    return png_paths

# PowerPoint COMオブジェクトを使ってスライドをPNGに変換する方法（Windows限定）
# Args:
#     pptx_path (str): 入力PPTXファイルパス
#     output_dir (str): 出力ディレクトリ
# Returns:
#     list[str]: 生成されたPNGファイルパスのリスト（スライド順）
def pptx_to_pngs_com(pptx_path, output_dir="slides_png"):
    import win32com.client
    import re

    # 出力ディレクトリを削除して再作成（初期化）
    if os.path.exists(output_dir):
        shutil.rmtree(output_dir)
    os.makedirs(output_dir)

    abs_pptx = str(Path(pptx_path).resolve())
    abs_out  = str(Path(output_dir).resolve())

    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = 1

    try:
        pres = ppt.Presentations.Open(abs_pptx)
        pres.Export(abs_out, "PNG")
        slide_count = pres.Slides.Count
        pres.Close()
    finally:
        ppt.Quit()

    # スライド順にソートしてパスリストを構築
    png_files = sorted(
        Path(abs_out).glob("*.PNG"),
        key=lambda p: int(re.search(r"\d+", p.stem).group())  # "スライド1", "スライド2" ... の順
    )

    # slide_001.png 形式にリネーム
    png_paths = []
    for i, src in enumerate(png_files, start=1):
        dst = Path(abs_out) / f"slide_{i:03d}.png"
        src.rename(dst)
        png_paths.append(str(dst))

    return png_paths

# PowerPoint COMオブジェクトが利用可能か確認する（Windows限定）
def is_powerpoint_available() -> bool:
    import platform
    if platform.system() != "Windows":
        return False
    try:
        import win32com.client
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Quit()
        return True
    except Exception:
        return False


# LibreOfficeを使ってスライドをPNGに変換する
# Args:
#     pptx_path (str): 入力PPTXファイルパス
#     output_dir (str): 出力ディレクトリ
# Returns:
#     list[str]: 生成されたPNGファイルパスのリスト（スライド順）
def pptx_to_pngs_libreoffice(pptx_path, output_dir="slides_png"):
    import subprocess

    if os.path.exists(output_dir):
        shutil.rmtree(output_dir)
    os.makedirs(output_dir)

    abs_pptx = str(Path(pptx_path).resolve())
    abs_out = str(Path(output_dir).resolve())

    # LibreOfficeのパス候補（環境に合わせて調整）
    libreoffice_candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",  # PATH が通っている場合
    ]
    libreoffice_path = None
    for candidate in libreoffice_candidates:
        if Path(candidate).exists() or candidate == "soffice":
            libreoffice_path = candidate
            break

    if libreoffice_path is None:
        raise FileNotFoundError("LibreOfficeが見つかりません。インストール先を確認してください。")

    subprocess.run(
        [
            libreoffice_path,
            "--headless",
            "--convert-to", "png",
            "--infilter=Impress PNG Export", 
            "--outdir", abs_out,
            abs_pptx,
        ],
        check=True,
    )

    # LibreOfficeの出力ファイル名は元ファイル名ベースになるため、
    # slide_001.png 形式にリネーム
    # 例: input.png, input2.png ... または input-001.png など環境依存
    png_files = sorted(Path(abs_out).glob("*.png"))
    png_paths = []
    for i, src in enumerate(png_files, start=1):
        dst = Path(abs_out) / f"slide_{i:03d}.png"
        src.rename(dst)
        png_paths.append(str(dst))

    return png_paths

# LibreOfficeが利用可能か確認
def is_libreoffice_available() -> bool:
    import platform
    import shutil
    # PATH上のsofficeのみで判定（クロスプラットフォーム対応）
    if shutil.which("soffice"):
        return True
    # Windowsのみパス候補を追加
    if platform.system() == "Windows":
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        return any(Path(c).exists() for c in candidates)
    return False

def _find_unoconv() -> str:
    found = shutil.which("unoconv")
    if found:
        return found
    raise FileNotFoundError("unoconv が見つかりません。sudo apt install unoconv でインストールしてください。")

# ──────────────────────────────────────────
# 2. スライドテキストから音声ファイルを生成
# ──────────────────────────────────────────
# 各スライドのテキストからMP3音声を生成します。
# Args:
#     pptx_path (str): 入力PPTXファイルパス
#     audio_dir (str): 音声ファイルの出力ディレクトリ
#     lang (str): 音声言語コード（デフォルト 'ja'）
#     voicevox (bool): VOICEVOX音声モード（デフォルト False）
#     voicevox=True の場合はローカルVOICEVOX APIを使用（WAV出力）。
#     voicevox=False の場合はgTTSを使用（MP3出力）。
# Returns:
#     list[str|None]: 各スライドの音声ファイルパス（空スライドはNone）
def generate_audio_files(pptx_path, audio_dir="slides_audio", lang="ja", voicevox=False, voicevoxid=3):

    # 出力ディレクトリを削除して再作成（初期化）
    if os.path.exists(audio_dir):
        shutil.rmtree(audio_dir)
    os.makedirs(audio_dir)

    prs = Presentation(pptx_path)
    audio_paths = []

    for i, slide in enumerate(prs.slides):
        # スライドノートがある場合はノートテキストを優先して取得
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        if notes_text:
            text = notes_text
            print(f"スライド {i+1}: ノートからテキスト取得", flush=True)
        else:
            # スライド上の全テキストを結合して取得
            text = " ".join(
                shape.text for shape in slide.shapes
                if hasattr(shape, "text")
            ).strip()
            print(f"スライド {i+1}: スライド本文からテキスト取得", flush=True)


        if text:
            audio_text_path = Path(audio_dir) / f"audio_{i+1:03d}.txt"
            with open(audio_text_path, "w", encoding="utf-8") as f:
                f.write(text)
                print(f"音声テキスト保存: {audio_text_path}", flush=True)

            # voicevoxがローカルで起動しているか確認
            if voicevox and is_voicevox_running():
                import requests
                VOICEVOX_URL = "http://localhost:50021"
                VOICEVOX_SPEAKER = voicevoxid  # 話者ID

                # VOICEVOX API で音声生成（WAV）
                try:
                    query = requests.post(
                        f"{VOICEVOX_URL}/audio_query",
                        params={"text": text, "speaker": VOICEVOX_SPEAKER}
                    ).json()

                    audio_bytes = requests.post(
                        f"{VOICEVOX_URL}/synthesis",
                        params={"speaker": VOICEVOX_SPEAKER},
                        json=query
                    ).content

                    audio_path = os.path.join(audio_dir, f"audio_{i+1:03d}.wav")
                    with open(audio_path, "wb") as f:
                        f.write(audio_bytes)
                    print(f"音声保存（VOICEVOX）: {audio_path}", flush=True)
                    audio_paths.append(audio_path)

                except Exception as e:
                    print(f"スライド {i+1}: VOICEVOX音声生成失敗 ({e}) → Noneとして処理", flush=True)
                    audio_paths.append(None)
            else:
                tts = gTTS(text=text, lang=lang, slow=False)
                audio_path = os.path.join(audio_dir, f"audio_{i+1:03d}.mp3")
                tts.save(audio_path)
                print(f"音声保存: {audio_path}", flush=True)
                audio_paths.append(audio_path)
        else:
            audio_paths.append(None)

    return audio_paths

# 空スライド用の無音WAVを生成（VOICEVOXモードで音声が生成できない場合のフォールバック）
def generate_silence_wav(output_path: str, duration_sec: float = 1.0) -> str:
    silence = AudioSegment.silent(duration=int(duration_sec * 1000))
    silence.export(output_path, format="wav")
    return output_path

# VOICEVOXがローカルで起動しているか確認する関数
def is_voicevox_running(url="http://localhost:50021"):
    import requests
    try:
        res = requests.get(f"{url}/version", timeout=2)
        return res.status_code == 200
    except requests.exceptions.ConnectionError:
        return False

# VOICEVOXの話者IDからクレジット表示文言を生成
# Args:
#     speaker_id (int): 使用する話者ID（style_id）
#     voicevox_url (str): VOICEVOX APIのベースURL
# Returns:
#     str: クレジット表示文言
def get_voicevox_credit(speaker_id: int, voicevox_url: str = "http://localhost:50021") -> str:
    import requests
    try:
        res = requests.get(f"{voicevox_url}/speakers", timeout=5)
        res.raise_for_status()
        speakers = res.json()
    except Exception as e:
        return f"VOICEVOX（クレジット情報取得失敗: {e}）"

    # speaker_id（style_id）からキャラクター名・スタイル名を特定
    for speaker in speakers:
        for style in speaker["styles"]:
            if style["id"] == speaker_id:
                character_name = speaker["name"]
                style_name = style["name"]
                return _build_credit(character_name, style_name)

    # speaker_id が style["id"] に直接対応する場合（上記でヒットしない場合のフォールバック）
    for speaker in speakers:
        for style in speaker["styles"]:
            if style["id"] == speaker_id:
                character_name = speaker["name"]
                style_name = style["name"]
                return _build_credit(character_name, style_name)

    return f"VOICEVOX（話者ID {speaker_id} が見つかりません）"

# キャラクター名・スタイル名から利用規約に応じたクレジット文言を生成
# VOICEVOX共通ルール:
#     - 基本形式: 「VOICEVOX:{キャラクター名}」
#     - 青山龍星: クレジット除去不可・明示必須
#     - その他: 同形式で対応
# Args:
#     character_name (str): キャラクター名（例: "ずんだもん"）
#     style_name (str): スタイル名（例: "ノーマル"）
# Returns:
#     str: クレジット表示文言
def _build_credit(character_name: str, style_name: str) -> str:
    # クレジット除去不可キャラ（規約上明示されているもの）
    NO_REMOVAL_CHARACTERS = {"青山龍星"}
    base_credit = f"VOICEVOX:{character_name}"
    if character_name in NO_REMOVAL_CHARACTERS:
        return f"{base_credit}（クレジット表記必須・除去不可）"
    return base_credit

def _find_font(size: int):
    from PIL import ImageFont
    candidates=[
        "C:/Windows/Fonts/meiryo.ttc",  # Windows
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",  # Ubuntu
        "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",  # macOS
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path,size)
        except Exception:
            continue
        return ImageFont.load_default()


# クレジット表示用のPNG画像を生成
# Args:
#     credit_text (str): 表示するクレジット文言（例: "VOICEVOX:ずんだもん"）
#     output_path (str): 出力PNGファイルパス
#     width (int): 画像の幅（px）
#     height (int): 画像の高さ（px）
#     bg_color (tuple): 背景色（RGB）
#     text_color (tuple): テキスト色（RGB）
#     font_size (int): フォントサイズ
# Returns:
#     str: 生成したPNGファイルパス
def generate_credit_slide(
    credit_text: str,
    output_path: str = "credit_slide.png",
    width: int = 1280,
    height: int = 720,
    bg_color: tuple = (0, 0, 0),
    text_color: tuple = (255, 255, 255),
    font_size: int = 40,
    image_path: str = None,
    image_max_size: tuple = (640, 480),
    margin: int = 20,
) -> str:
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (width, height), color=bg_color)
    draw = ImageDraw.Draw(img)

    # フォント読み込み（システムフォントにフォールバック）
    font = _find_font(font_size)

    # 画像を中央配置
    if image_path and os.path.exists(image_path):
        overlay = Image.open(image_path).convert("RGBA")
        ow, oh = overlay.size
        max_w, max_h = image_max_size
        scale = min(max_w / ow, max_h / oh, 1.0)
        new_w = int(ow * scale)
        new_h = int(oh * scale)
        overlay = overlay.resize((new_w, new_h), Image.LANCZOS)

        # 画像を画面中央に配置
        x = (width - new_w) // 2
        y = (height - new_h) // 2
        img.paste(overlay, (x, y), overlay)

    # テキストの位置を計算して描画
    bbox = draw.textbbox((0, 0), credit_text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    text_x = width - text_width - margin
    text_y = height - text_height - margin

    draw.text((text_x, text_y), credit_text, fill=text_color, font=font)
    img.save(output_path, "PNG")
    print(f"クレジットスライド生成: {output_path}", flush=True)
    return output_path

# ──────────────────────────────────────────
# 3. 音声ファイルを結合してWAVに出力
# ──────────────────────────────────────────
# 複数のMP3音声ファイルを結合してWAVファイルを生成します。
# Args:
#     audio_paths (list[str|None]): 音声ファイルパスのリスト
#     output_path (str): 結合後のWAVファイルパス
# Returns:
#     str: 出力WAVファイルパス
def combine_audio(audio_paths, output_path="combined_audio.wav"):
    combined = AudioSegment.empty()
    for path in audio_paths:
        if path and os.path.exists(path):
            combined += AudioSegment.from_mp3(path)
        else:
            # 空スライドは2秒の無音を挿入
            combined += AudioSegment.silent(duration=2000)

    combined.export(output_path, format="wav")
    print(f"音声結合完了: {output_path}", flush=True)
    return output_path


# ──────────────────────────────────────────
# 4. FFmpegでPNG + 音声を動画に合成
# ──────────────────────────────────────────
# PNG画像リストと音声ファイルをFFmpegで動画に合成します。
# 各スライドの表示時間は対応する音声の長さに合わせます。
# Args:
#     png_paths (list[str]): PNGファイルパスのリスト（スライド順）
#     audio_paths (list[str|None]): 音声ファイルパスのリスト
#     output_mp4 (str): 出力MP4ファイルパス
#     debug (bool): デバッグモード
#     quality (int): 動画品質（1-31、値が小さいほど高品質）
def create_video_ffmpeg(png_paths, audio_paths, output_mp4="output.mp4", debug=False, quality=5):
    import subprocess
    # FFmpegのパスを取得（同ディレクトリのffmpeg.exeを想定）
    _BASE_DIR = Path(__file__).parent
    ffmpeg_path = find_ffmpeg("ffmpeg")

    # 各スライドの音声長を取得
    durations = []
    for path in audio_paths:
        if path and os.path.exists(path):
            seg = AudioSegment.from_mp3(path)
            durations.append(len(seg) / 1000.0)
        else:
            durations.append(2.0)  # 空スライドは2秒
    total_duration = sum(durations)  # durationリストの合計
    print(f"合計再生時間: {total_duration:.3f} 秒", flush=True)

    # concat demuxerファイルを生成（スライドごとに個別duration）
    concat_file = "concat_list.txt"
    with open(concat_file, "w", encoding="utf-8") as f:
        for png, duration in zip(png_paths, durations):
            f.write(f"file '{os.path.abspath(png)}'\n")
            f.write(f"duration {duration:.3f}\n")
        # 最後のフレームを明示（FFmpegのconcat demuxer仕様）
        f.write(f"file '{os.path.abspath(png_paths[-1])}'\n")

    # 音声を結合してWAVに出力
    combined_audio_path = "combined_audio.wav"
    combined_audio = combine_audio(audio_paths, output_path=combined_audio_path)

    print(f"結合された音声の長さ: {len(AudioSegment.from_wav(combined_audio_path)) / 1000.0} 秒", flush=True)

    cmd = [
        ffmpeg_path,                # 実行するffmpegのパス（スクリプトと同階層のffmpeg.exe）
        "-y",                       # 出力ファイルが既に存在する場合、確認なしで上書き
        "-f", "concat",             # 入力フォーマットとしてconcat demuxerを使用
        "-safe", "0",               # concat_list.txt内の絶対パスを許可（デフォルトは相対パスのみ）
        "-i", concat_file,          # 入力①：concat_list.txt（スライドPNGと各表示時間を定義したファイル）
        "-i", combined_audio,       # 入力②：結合済み音声ファイル（WAV）
        "-c:v", "mpeg4",            # 映像コーデックにMPEG-4 Part 2（FFmpeg完全内蔵）
        "-q:v", str(quality),       # 品質指定（1〜31、値が小さいほど高品質）
        "-pix_fmt", "yuv420p",      # ピクセルフォーマットをYUV 4:2:0に変換（広い互換性のため）
        "-vf", "format=rgb24,scale=trunc(iw/2)*2:trunc(ih/2)*2", # RGBビット深度を24にし、幅・高さをそれぞれ2の倍数に切り捨てる
        "-r", "30",                 # フレームレートを30fpsに設定
        "-t", str(total_duration),  # 動画の総再生時間を指定（concat_list.txtのduration合計）
        "-c:a", "aac",              # 音声コーデックにAACを使用
        "-map", "0:v:0",            # 入力①（concat_list.txt）の映像ストリーム0番を出力に使用
        "-map", "1:a:0",            # 入力②（combined_audio）の音声ストリーム0番を出力に使用
        output_mp4                  # 出力ファイルパス
    ]

    print(f"FFmpegコマンド: {' '.join(cmd)}", flush=True)

    # FFmpegの出力をリアルタイムで表示しながら実行
    proc = subprocess.Popen(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        bufsize=1,
    )
    for raw_line in proc.stdout:
        line = raw_line.decode("utf-8", errors="replace")
        sys.stdout.write(line)
        sys.stdout.flush()

    # FFmpegのプロセスが終了するまで待機し、終了コードを確認
    proc.wait()
    if proc.returncode != 0:
        raise subprocess.CalledProcessError(proc.returncode, cmd)

    print(f"動画生成完了: {output_mp4}", flush=True)

    # デバッグモードのみ中間ファイルを保持
    if not debug:
        os.remove(concat_file)
        os.remove(combined_audio)
        print("中間ファイルを削除しました。", flush=True)
    else:
        print(f"[DEBUG] 中間ファイルを保持しています: {concat_file}, {combined_audio}", flush=True)


# ──────────────────────────────────────────
# 5. メイン処理（統合実行）
# ──────────────────────────────────────────
# PPTXファイルを動画（MP4）に変換します。
# LibreOfficeを使わず、python-pptx + Pillow + gTTS + FFmpegで実装。

# Args:
#     pptx_path (str): 入力PPTXファイルパス
#     output_mp4 (str): 出力MP4ファイルパス
#     dpi (int): PNG解像度
#     lang (str): 音声言語コード
#     png_dir (str): PNG一時保存ディレクトリ
#     audio_dir (str): 音声一時保存ディレクトリ
#     voicevox (bool): VOICEVOX音声モード
#     voicevoxid (int): VOICEVOX話者ID(デフォルト3: ずんだもんノーマル)
#     creditimg (str): クレジット画像ファイルパス（例: /file/to/credit.png）
#     debug (bool): デバッグモード（中間ファイルを保持）
def pptx_to_video(
    pptx_path,
    output_mp4="output.mp4",
    dpi=150,
    quality=5,
    lang="ja",
    png_dir="slides_png",
    audio_dir="slides_audio",
    voicevox=False,
    voicevoxid=3,
    creditimg=None,
    creditbg=None,
    creditcolor=None,
    debug=False,
):

    # ── 環境判定（最初に1回だけ実施） ──────────────────────
    import platform
    ENV_OS              = platform.system()
    ENV_USE_PPT         = is_powerpoint_available()   # Windows + PowerPoint
    ENV_USE_LIBREOFFICE = is_libreoffice_available()  # LibreOffice
    ENV_USE_VOICEVOX    = is_voicevox_running()        # VOICEVOX APIサーバー

    print(f"[ENV] OS={ENV_OS}, PowerPoint={ENV_USE_PPT}, LibreOffice={ENV_USE_LIBREOFFICE}, VOICEVOX={ENV_USE_VOICEVOX}", flush=True)

    print("=== STEP 1: PNG変換 ===", flush=True)
    if ENV_USE_PPT:
        print("PowerPoint COMを使用してPNG変換します。", flush=True)
        png_paths = pptx_to_pngs_com(pptx_path, output_dir=png_dir)
    elif ENV_USE_LIBREOFFICE:
        print("LibreOfficeを使用してPNG変換します。", flush=True)
        png_paths = pptx_to_pngs_libreoffice(pptx_path, output_dir=png_dir)
    else:
        print("PowerPoint・LibreOfficeが見つかりません。python-pptx + Pillowでフォールバック変換します。", flush=True)
        png_paths = pptx_to_pngs(pptx_path, output_dir=png_dir, dpi=dpi)

    # クレジットスライドを末尾に追加
    if voicevox and is_voicevox_running():
        speaker_id = voicevoxid  # 実際に使用している話者ID
        credit_text = get_voicevox_credit(speaker_id)
    else:
        credit_text = ''

    if creditimg and os.path.exists(creditimg):
        creditimg_path=creditimg
    else:
        creditimg_path = None

    if creditbg != '' and creditbg is not None:
        # カラーコードをRGBタプルに変換
        creditbg_r, creditbg_g, creditbg_b = bytes.fromhex(creditbg.lstrip('#'))
        creditbg_value = (creditbg_r, creditbg_g, creditbg_b)
    else:
        creditbg_value = (255, 255, 255)

    if creditcolor != '' and creditcolor is not None:
        # カラーコードをRGBタプルに変換
        creditcolor_r, creditcolor_g, creditcolor_b = bytes.fromhex(creditcolor.lstrip('#'))
        creditcolor_value = (creditcolor_r, creditcolor_g, creditcolor_b)
    else:
        creditcolor_value = (128, 128, 128)

    if quality != '' and quality is not None:
        if isinstance(quality, int) and 1 <= quality <= 31:
            quality_value = quality
        else:
            quality_value = 5  # デフォルト品質
    else:
        quality_value = 5  # デフォルト品質

    if credit_text != '' or (creditimg_path and os.path.exists(creditimg_path)):
        credit_png = generate_credit_slide(
            credit_text=credit_text,
            output_path=os.path.join(os.path.abspath(png_dir), "slide_credit.png"),
            image_path=creditimg_path,
            bg_color=creditbg_value,
            text_color=creditcolor_value,
        )
        png_paths.append(credit_png)
    else:
        print("クレジットスライドは生成されませんでした（テキスト・画像ともに指定なし）。", flush=True)

    print("\n=== STEP 2: 音声生成 ===", flush=True)
    audio_paths = generate_audio_files(pptx_path, audio_dir=audio_dir, lang=lang, voicevox=voicevox, voicevoxid=voicevoxid)

    # クレジット用無音を末尾に追加
    if voicevox and is_voicevox_running():
        silence_path = os.path.join(os.path.abspath(audio_dir), "audio_credit.wav")
        audio_paths.append(generate_silence_wav(silence_path, duration_sec=1.0))

    print("\n=== STEP 3: 動画合成 ===", flush=True)
    create_video_ffmpeg(png_paths, audio_paths, output_mp4=output_mp4, debug=debug, quality=quality_value)

    # 各スライドの音声ファイルを削除（デバッグ時は保持）
    if not debug:
        for p in png_paths:
            if p and os.path.exists(p):
                os.remove(p)
        print("画像ファイルを削除しました。", flush=True)
        for p in audio_paths:
            if p and os.path.exists(p):
                os.remove(p)
        print("音声ファイルを削除しました。", flush=True)
    else:
        print("[DEBUG] 画像ファイルを保持しています。", flush=True)
        print("[DEBUG] 音声ファイルを保持しています。", flush=True)

    print(f"\n✅ 完了: {output_mp4}", flush=True)

# #RRGGBB 形式の文字列を (R, G, B) タプルに変換
def hex_color(value: str) -> tuple[int, int, int]:
    import re
    m = re.fullmatch(r'#([0-9A-Fa-f]{2})([0-9A-Fa-f]{2})([0-9A-Fa-f]{2})', value)
    if not m:
        raise argparse.ArgumentTypeError(f'カラーコードの形式が不正です: {value}（例: #FF00FF）')
    return tuple(int(m.group(i), 16) for i in range(1, 4))


# ──────────────────────────────────────────
# 実行
# ──────────────────────────────────────────
def main():
    args = doArgParse()
    
    # デバッグモードの場合、標準出力をファイルにも保存するためのクラス
    if args['debug']:
        import sys
        log_path = os.path.splitext(os.path.abspath(args['output']))[0] + "_debug.log"
        _tee = _Tee(sys.stdout, log_path)  # reconfigure もここで完結
        sys.stdout = _tee

    print('====== Slide to Movie ======', flush=True)
    print('                     v.0.0.12', flush=True)
    print(f'指定された引数: {args}', flush=True)
        
    pptx_to_video(
        pptx_path=args['file'],
        output_mp4=args['output'],
        dpi=args['dpi'],
        quality=args['quality'],
        lang=args['lang'],
        voicevox=args['voicevox'],
        voicevoxid=args['voicevoxid'],
        creditimg=args['creditimg'],
        creditbg=args['creditbg'],
        creditcolor=args['creditcolor'],
        debug=args['debug'],
    )

# 相対パス取得
def getRelativePath(filePath):
    from pathlib import Path
    if(Path(filePath).is_absolute()):
        return Path(filePath).relative_to(Path.cwd())
    else:
        return filePath

# 絶対パス取得
def getAbsolutePath(filePath):
    from pathlib import Path
    if(Path(filePath).is_absolute()):
        return filePath
    else:
        return Path(filePath).resolve()

# 引数解析
def doArgParse() -> dict:
    parser = build_parser(ARG_DESCRIPTION)
    args = parser.parse_args()
    arg_dict = vars(args)
    config_path = arg_dict.pop('config')
    ini_dict = load_ini(config_path)

    # 優先順位: 引数 > ini > デフォルト
    merged = {}
    for opt in OPTION_DEFS:
        key = opt['name']
        if arg_dict.get(key) is not None:
            merged[key] = arg_dict[key]
        elif key in ini_dict:
            merged[key] = ini_dict[key]
        else:
            merged[key] = opt['default']

    # 必須チェック
    for opt in OPTION_DEFS:
        if opt['required'] and not merged.get(opt['name']):
            parser.error(f'--{opt["name"]} が指定されていません（引数またはiniファイルで設定してください）')

    # iniファイルが存在しない場合は現在の設定値で生成する
    if Path(config_path).exists()==False:
        save_ini(config_path, merged)

    return merged


# ──────────────
# オプション定義
# ──────────────
from external_define import OPTION_DEFS, ARG_DESCRIPTION, CONFIG_DEFAULT

# ──────────────
# INIファイルの読み書き関数
# ──────────────
from external_define import load_ini, save_ini

# OPTION_DEFSからargparseを生成
def build_parser(desc=ARG_DESCRIPTION) -> argparse.ArgumentParser:
    import argparse
    parser = argparse.ArgumentParser(description=desc)
    parser.add_argument('--config', default=CONFIG_DEFAULT, help='設定ファイルパス')

    for opt in OPTION_DEFS:
        if opt['store_true']:
            parser.add_argument(f'--{opt["name"]}', action=argparse.BooleanOptionalAction, default=None, help=opt['help'])
        else:
            parser.add_argument(f'--{opt["name"]}', type=opt['type'], default=None, help=opt['help'])

    return parser


if __name__ == '__main__':
    main()